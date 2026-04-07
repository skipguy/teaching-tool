"""
ppt_gen.py — Generate a PowerPoint file with one question per slide.

Each slide contains:
  - The cropped question image, centred and scaled to fill the slide.
  - A slide number label in the bottom-right corner.
"""

from __future__ import annotations

import io
import logging
from typing import Sequence

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

logger = logging.getLogger(__name__)

# Slide dimensions  (16:9)
SLIDE_W_CM = 25.4
SLIDE_H_CM = 14.29
SLIDE_W_EMU = Cm(SLIDE_W_CM)
SLIDE_H_EMU = Cm(SLIDE_H_CM)

# Margins inside the slide for the question image (in cm)
MARGIN_CM = 0.5
MARGIN_EMU = Cm(MARGIN_CM)

TOPIC_LABEL_H_EMU = Cm(1.0)   # reserved height at top for topic label

MAX_IMG_W_EMU = SLIDE_W_EMU - MARGIN_EMU * 2
# leave room for topic label at top and slide number at bottom
MAX_IMG_H_EMU = SLIDE_H_EMU - MARGIN_EMU * 2 - Cm(0.8) - TOPIC_LABEL_H_EMU

SLIDE_NUM_FONT_SIZE = Pt(10)
SLIDE_NUM_COLOR = RGBColor(0x99, 0x99, 0x99)

TOPIC_FONT_SIZE = Pt(28)
TOPIC_COLOR = RGBColor(0xCC, 0x00, 0x00)   # red


def generate_ppt(
    question_images: list[bytes],
    question_numbers: list[int | str],
    topic_labels: list[str] | None = None,  # one label per slide (or None)
) -> bytes:
    """
    Build a pptx in memory and return its bytes.

    Parameters
    ----------
    question_images  : PNG bytes per question, already in desired order
    question_numbers : display labels for each question (same order)
    topic_labels     : optional topic string per slide shown top-left in red
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU

    blank_layout = prs.slide_layouts[6]

    for slide_idx, (img_bytes, q_num) in enumerate(
        zip(question_images, question_numbers), start=1
    ):
        slide = prs.slides.add_slide(blank_layout)

        # ── background: white ────────────────────────────────────────────────
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # ── topic label top-left (red) ───────────────────────────────────────
        label_text = (topic_labels[slide_idx - 1] if topic_labels else "").strip()
        if label_text:
            lbl_w = SLIDE_W_EMU - Cm(1.0)
            lbl_h = TOPIC_LABEL_H_EMU
            txBox = slide.shapes.add_textbox(Cm(0.5), Cm(0.15), lbl_w, lbl_h)
            tf = txBox.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = label_text
            run.font.size = TOPIC_FONT_SIZE
            run.font.bold = True
            run.font.color.rgb = TOPIC_COLOR

        # ── insert question image ────────────────────────────────────────────
        img_stream = io.BytesIO(img_bytes)
        pil_img = Image.open(img_stream)
        img_w_px, img_h_px = pil_img.size

        scale_w = MAX_IMG_W_EMU / _px_to_emu(img_w_px)
        scale_h = MAX_IMG_H_EMU / _px_to_emu(img_h_px)
        scale = min(scale_w, scale_h, 1.0)

        emu_w = int(_px_to_emu(img_w_px) * scale)
        emu_h = int(_px_to_emu(img_h_px) * scale)

        # Centre horizontally; vertically within the area below the topic label
        img_area_top = TOPIC_LABEL_H_EMU + MARGIN_EMU
        img_area_h   = SLIDE_H_EMU - img_area_top - Cm(0.8)
        left = (SLIDE_W_EMU - emu_w) // 2
        top  = img_area_top + (img_area_h - emu_h) // 2

        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, left, top, width=emu_w, height=emu_h)

        # ── slide number bottom-right ─────────────────────────────────────────
        num_w = Cm(3.0)
        num_h = Cm(0.7)
        num_left = SLIDE_W_EMU - num_w - Cm(0.3)
        num_top  = SLIDE_H_EMU - num_h - Cm(0.2)
        txBox = slide.shapes.add_textbox(num_left, num_top, num_w, num_h)
        tf = txBox.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.RIGHT
        run = para.add_run()
        run.text = f"{slide_idx}"
        run.font.size = SLIDE_NUM_FONT_SIZE
        run.font.color.rgb = SLIDE_NUM_COLOR

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# EMU per pixel at 96 dpi (PowerPoint default screen resolution)
# 1 inch = 914400 EMU, 1 inch = 96 px  →  9525 EMU/px
_EMU_PER_PX = 914400 / 96


def _px_to_emu(px: int) -> int:
    return int(px * _EMU_PER_PX)
